#!/usr/bin/env python3
"""Generate Hassan_Hackathon_Deck.pptx for MOEI × 42AD AgentEra elimination round."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
INK   = RGBColor(0x0C, 0x1B, 0x2A)
RED   = RGBColor(0xEF, 0x33, 0x40)
GREEN = RGBColor(0x00, 0x96, 0x39)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF2, 0xF5)
SLATE = RGBColor(0xB0, 0xC4, 0xDE)
DGRAY = RGBColor(0x40, 0x50, 0x60)
PANEL = RGBColor(0x14, 0x2A, 0x40)
DARK2 = RGBColor(0x06, 0x10, 0x1A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(blank)


def box(sl, x, y, w, h, fill, line_color=None):
    from pptx.util import Inches
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp


def txt(sl, text, x, y, w, h, size=18, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def notes(sl, text):
    ns = sl.notes_slide
    ns.notes_text_frame.text = text


# ── SLIDE 1 · Title ───────────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.3, 7.5, RED)
box(s, 0.3, 7.42, 13.03, 0.08, GREEN)

txt(s, "HASSAN", 0.5, 0.7, 12.3, 2.2, size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Hyper-Adaptive Smart Service Agent for the Nation",
    0.5, 2.85, 12.3, 0.8, size=22, color=SLATE, align=PP_ALIGN.CENTER)
box(s, 3.5, 3.8, 6.33, 0.05, RED)
txt(s, "MOEI  x  42 Abu Dhabi  |  AgentEra Hackathon  |  Challenge 3",
    0.5, 3.95, 12.3, 0.6, size=17, color=SLATE, align=PP_ALIGN.CENTER)
txt(s, "Team Hassan",
    0.5, 4.75, 12.3, 0.7, size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Arabic subtitle
txt(s, "خدمة حكومة الإمارات بذكاء يستحقه المواطن",
    0.5, 5.7, 12.3, 0.7, size=18, color=RGBColor(0x70, 0x90, 0xB0), align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A opens (30 seconds max):
"Good morning. In the next five minutes, we're going to show you something that doesn't exist in any UAE government service today: a single AI agent that knows a citizen across every channel, never makes them repeat themselves, and acts before they even complain. This is Hassan."

[Pause one second. Let the name land. Move to slide 2.]""")


# ── SLIDE 2 · Citizen Story ───────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, RED)
box(s, 0.15, 0, 13.18, 0.95, INK)
txt(s, "Meet Fatima Al Mansouri  —  Ajman  —  SZHP Housing Case",
    0.4, 0.1, 12.5, 0.75, size=26, bold=True, color=WHITE)

# TODAY column
box(s, 0.3, 1.05, 5.9, 6.1, RGBColor(0xFD, 0xEC, 0xEC))
txt(s, "TODAY", 0.55, 1.12, 5.5, 0.45, size=13, bold=True, color=RED)
rows_today = [
    "4 months behind on SZHP payments after a medical emergency",
    "Calls MOEI hotline — 40-minute wait, transferred twice",
    "Re-explains herself from scratch on every channel",
    "Tries WhatsApp — no memory of the call",
    "5 DAYS  |  4 CHANNELS  |  1 HOUR WAIT  |  No answer",
]
for i, row in enumerate(rows_today):
    bold_row = i == 4
    col = RED if i == 4 else DGRAY
    size = 15 if i == 4 else 14
    txt(s, row, 0.55, 1.62 + i * 1.05, 5.5, 0.9, size=size, bold=bold_row, color=col)

# Arrow
txt(s, ">", 6.3, 3.5, 0.75, 0.8, size=44, bold=True, color=INK, align=PP_ALIGN.CENTER)

# WITH HASSAN column
box(s, 7.15, 1.05, 5.9, 6.1, RGBColor(0xE8, 0xF8, 0xED))
txt(s, "WITH HASSAN", 7.35, 1.12, 5.5, 0.45, size=13, bold=True, color=GREEN)
rows_hassan = [
    "Sends one WhatsApp — Hassan already knows her",
    "SZHP rules engine: rescheduling plan in 8 seconds",
    "Case created automatically. SLA clock starts.",
    "Proactive status update sent before she has to ask",
    "90 SECONDS  |  1 CONVERSATION  |  Never repeats",
]
for i, row in enumerate(rows_hassan):
    bold_row = i == 4
    col = GREEN if i == 4 else DGRAY
    size = 15 if i == 4 else 14
    txt(s, row, 7.35, 1.62 + i * 1.05, 5.5, 0.9, size=size, bold=bold_row, color=col)

notes(s, """PRESENTER A (storytelling — slow and vivid):
"Fatima Al Mansouri. She lives in Ajman. Four months behind on her Sheikh Zayed Housing Programme payments after a medical emergency."

[Point left]
"Today: five days, four channels, one hour waiting, re-explaining herself every single time. No resolution."

[Sweep right]
"With Hassan: she sends one WhatsApp. Ninety seconds later the system knows who she is, runs the housing rules engine, proposes a rescheduling plan, creates a case, and sends her a proactive update before she even has to ask again."

"One conversation. Every channel. She never repeats herself once."

[Move to slide 3 — fast.]""")


# ── SLIDE 3 · The Federal Problem ─────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.15, 7.5, GREEN)

txt(s, "The Problem Is Bigger Than Fatima",
    0.45, 0.15, 12.5, 0.85, size=32, bold=True, color=WHITE)
box(s, 0.45, 1.05, 12.43, 0.05, RED)

# 3 stat blocks
stats = [
    ("174", "federal services\nat MOEI"),
    ("4", "channels per service\nweb  ·  voice  ·  WhatsApp  ·  mobile"),
    ("0", "shared citizen context\nbetween any of them — today"),
]
for i, (num, label) in enumerate(stats):
    x = 0.55 + i * 4.2
    box(s, x, 1.2, 3.7, 2.6, PANEL)
    txt(s, num, x+0.1, 1.28, 3.5, 1.5, size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.1, 2.72, 3.5, 0.95, size=15, color=SLATE, align=PP_ALIGN.CENTER)

txt(s, "Siloed CRMs  ·  Manual handoffs  ·  No cross-channel memory  ·  Reactive, not proactive",
    0.4, 3.95, 12.5, 0.6, size=17, color=SLATE, align=PP_ALIGN.CENTER)

box(s, 0.4, 4.65, 12.53, 0.05, GREEN)
txt(s, "UAE National AI Strategy 2031  ·  AED 335 billion digital economy target",
    0.4, 4.8, 12.5, 0.55, size=16, color=GREEN, align=PP_ALIGN.CENTER)

box(s, 1.0, 5.55, 11.33, 1.25, RGBColor(0x18, 0x32, 0x4C))
txt(s, "No team here is solving the cross-channel handoff.  We are.",
    1.1, 5.65, 11.1, 1.05, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A (confident, fast):
"This isn't just Fatima's problem. MOEI manages 174 federal services across four channels — and right now, none of them share context. Citizens re-explain themselves every single time. Agents work blind."

"The UAE National AI Strategy 2031 calls for exactly this kind of transformation — and nobody in this room has solved the cross-channel handoff."

[Beat.]

"We have. And we're going to prove it right now."

[Slide 4 — architecture, 20 seconds.]""")


# ── SLIDE 4 · Architecture ────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, INK)
box(s, 0.15, 0, 13.18, 0.9, INK)
txt(s, "One Brain. Every Channel.", 0.45, 0.1, 12.0, 0.75, size=30, bold=True, color=WHITE)

# Channel row
channels = ["WhatsApp", "Voice", "Web / Chat", "Mobile"]
for i, ch in enumerate(channels):
    x = 0.8 + i * 2.95
    box(s, x, 1.0, 2.5, 0.75, INK)
    txt(s, ch, x+0.05, 1.08, 2.4, 0.58, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "v v v v", 0.8, 1.8, 11.5, 0.45, size=18, color=INK, align=PP_ALIGN.CENTER)

# Gateway
box(s, 1.1, 2.28, 11.13, 0.78, PANEL)
txt(s, "FastAPI Channel Gateway  |  correlation-ID  |  language detection  |  PII redaction at entry",
    1.2, 2.38, 10.9, 0.58, size=15, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "v", 6.4, 3.1, 0.7, 0.5, size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)

# Supervisor
box(s, 1.1, 3.62, 11.13, 0.82, RED)
txt(s, "LangGraph Supervisor  |  Router > Sentiment > Memory > Guardrails > Dispatcher > Escalation > Composer",
    1.2, 3.72, 10.9, 0.62, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "v", 6.4, 4.47, 0.7, 0.5, size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)

# Worker agents
workers = ["Housing\n(SZHP)", "Energy\n(Tariffs)", "Transport &\nMaritime", "Knowledge\n(FTS / RAG)", "CRM\n(Profile)"]
for i, w in enumerate(workers):
    x = 0.5 + i * 2.5
    box(s, x, 4.98, 2.2, 1.05, RGBColor(0x00, 0x70, 0x29))
    txt(s, w, x+0.05, 5.06, 2.1, 0.9, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data layer
txt(s, "PostgreSQL (system of record)  |  Redis (cross-channel session)  |  Langfuse (tracing + audit)  |  Postgres FTS (bilingual RAG)",
    0.4, 6.18, 12.5, 0.5, size=12, color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A (20 seconds — fast architecture overview):
"One LangGraph supervisor. Every channel feeds into it. It routes to specialist agents — housing, energy, transport, maritime. Cross-channel memory in Redis and Postgres. The citizen's profile is unified."

"Now let's see it live."

[Hand off to PRESENTER B on the laptop. Move to slide 5 and leave it on screen.]""")


# ── SLIDE 5 · Demo Cue ────────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.3, 7.5, RED)
box(s, 0.3, 7.42, 13.03, 0.08, GREEN)

txt(s, "LIVE DEMO", 0.5, 1.2, 12.3, 2.3, size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Fatima's journey  —  four channels  —  one memory  —  in real time.",
    0.5, 3.55, 12.3, 0.75, size=22, color=SLATE, align=PP_ALIGN.CENTER)
box(s, 2.0, 4.45, 9.33, 0.06, RED)
txt(s, "Web Chat   >   WhatsApp   >   Proactive Update   >   Human Escalation",
    0.5, 4.6, 12.3, 0.65, size=20, color=GREEN, align=PP_ALIGN.CENTER)

notes(s, """[LEAVE THIS SLIDE ON SCREEN DURING THE ENTIRE DEMO — ~2.5 minutes]

PRESENTER B runs the demo. PRESENTER A narrates each step out loud:

STEP 1 — Web Chat (localhost:3000/chat or ngrok URL)
  - Sign in as Fatima Al Mansouri (UAE PASS mock, ID: 784199011810004)
  - Type: "I'm behind on my Sheikh Zayed housing payments after a medical emergency — what are my options?"
  NARRATE: "The SZHP rules engine runs in real time — deterministic policy, not an LLM guess. A case is created."

STEP 2 — WhatsApp
  - Send the same citizen's WhatsApp message: "What's the status of my request?"
  NARRATE: "Different channel. Hassan already knows her — no re-introduction. It pulls her unified profile by Customer ID and returns the real case status with SLA timing."

STEP 3 — Admin > Citizens > Fatima > Send Status Update
  NARRATE: "Her case is SLA-breached. We send a proactive update before she has to chase. That is the difference."

STEP 4 — Admin > Agent Co-pilot
  NARRATE: "She's flagged as a Repeat Escalator. Two dataset signals fired — reopened case plus repeat history. Hassan routes her to a human automatically. The agent sees her entire cross-channel history on one screen. Acts in one click."

PRESENTER C (quietly, while demo runs):
"Note: voice and mobile run through the exact same supervisor. Same memory, same profile."

TARGET: finish demo by 3:30. DO NOT RUSH STEP 2 — that's the money shot.""")


# ── SLIDE 6 · Agentic Depth ───────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, RED)
box(s, 0.15, 0, 13.18, 0.9, INK)
txt(s, "Agentic Depth", 0.45, 0.1, 9.5, 0.75, size=30, bold=True, color=WHITE)
txt(s, "25 pts", 10.2, 0.1, 2.8, 0.75, size=30, bold=True, color=RED, align=PP_ALIGN.RIGHT)

# Node pipeline
nodes = [
    ("Router", INK), ("Sentiment", INK), ("Memory", INK),
    ("Guardrails", RGBColor(0x70, 0x00, 0x10)),
    ("Dispatcher", RED), ("Escalation", RED),
    ("Composer", INK), ("Persist", RGBColor(0x00, 0x70, 0x29)),
]
for i, (name, col) in enumerate(nodes):
    x = 0.3 + i * 1.6
    box(s, x, 1.0, 1.35, 0.65, col)
    txt(s, name, x+0.02, 1.08, 1.31, 0.5, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        txt(s, ">", x + 1.35, 1.12, 0.25, 0.35, size=13, color=INK, align=PP_ALIGN.CENTER)

# 4 capability blocks
caps = [
    ("Multi-agent specialist dispatch",
     "Housing · Energy · Transport · Maritime · Infrastructure — each with dedicated tools and deterministic sub-logic."),
    ("Predictive escalation engine",
     "6 CRM signals fused (anger · SLA breach · reopened · repeat escalator · critical · VIP). Escalates on >= 2. Predicts before citizen complains."),
    ("Deterministic policy engine",
     "SZHP rules in Python, not LLM. Every decision is citable: 'SZHP-R3.1: arrears >= 3 months, eligible for hardship rescheduling.' Right-to-explanation built in."),
    ("Full PDPL audit trail",
     "Every node input/output captured to Langfuse. Every decision persisted to Postgres. Citizen-accessible. Replayable for any auditor."),
]
for i, (title, desc) in enumerate(caps):
    y = 1.88 + i * 1.27
    box(s, 0.3, y, 12.73, 1.18, LGRAY)
    box(s, 0.3, y, 0.08, 1.18, RED if i % 2 == 0 else GREEN)
    txt(s, title, 0.55, y + 0.1, 12.0, 0.45, size=16, bold=True, color=INK)
    txt(s, desc,   0.55, y + 0.57, 12.0, 0.55, size=13, color=DGRAY)

notes(s, """[BACKUP SLIDE — show during Q&A if asked "how does the agent actually work?"]

Key points:
- Router uses fast 8B model; Reasoner uses Llama 3.3 70B — role-based routing, not hardcoded models
- Escalation node is dataset-driven, not LLM confidence — it fires on measurable CRM signals
- Rules engine means every housing decision is explainable, not a black box
- Langfuse captures every node — full audit trail for PDPL Art. 7 right-to-explanation""")


# ── SLIDE 7 · Federal-Grade ───────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.15, 7.5, GREEN)
box(s, 0.15, 0, 13.18, 0.9, DARK2)
txt(s, "Federal-Grade. Built for Production.", 0.45, 0.1, 10.0, 0.75, size=30, bold=True, color=WHITE)
txt(s, "25 pts", 10.2, 0.1, 2.8, 0.75, size=30, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

federal = [
    ("PDPL Compliant", "Art. 4 Consent  |  Art. 5 Purpose limitation  |  Art. 6 UAE residency  |  Art. 7 Access  |  Art. 8 Security  |  Art. 9 Breach SOP"),
    ("UAE PASS Identity", "Synthetic identities in demo; real UAE PASS sandbox in Phase 1 pilot. Verified citizen flow — no guessing who Fatima is."),
    ("Data Residency", "Azure UAE North / Core42 (G42 Cloud). Fully containerised — on-prem deployable if TDRA requires. Data never leaves UAE."),
    ("PII + Bias Guards", "Emirates ID / IBAN / mobile / email redacted at gateway AND composer. Nationality / gender / religion bias detector built in."),
    ("Audit Trail UI", "Every decision clickable by the citizen or any auditor. PDPL Art. 7 right-to-explanation: exact rule citations, not LLM post-hoc rationalisation."),
    ("Provider Failover", "Groq > Claude/OpenAI > Gemini. Automatic cascade. No single-vendor lock-in. Zero hard-fails for the citizen."),
]
for i, (title, desc) in enumerate(federal):
    col = 0 if i < 3 else 1
    row = i % 3
    x = 0.4 + col * 6.5
    y = 1.05 + row * 1.9
    box(s, x, y, 6.2, 1.75, PANEL)
    box(s, x, y, 6.2, 0.1, GREEN if col == 0 else RED)
    txt(s, title, x+0.2, y+0.18, 5.8, 0.45, size=16, bold=True, color=GREEN)
    txt(s, desc,  x+0.2, y+0.65, 5.8, 1.0, size=12, color=SLATE)

notes(s, """[BACKUP SLIDE — use for Q&A on compliance / data privacy / production readiness]

Key Q&A answers:
Q: "What about data privacy?" -> PDPL article-by-article mapping, PII redacted at two independent layers, full audit trail
Q: "Where does the data live?" -> Azure UAE North or Core42, containerised, on-prem option available
Q: "What if an LLM provider goes down?" -> Automatic cascade: Groq -> Claude/OpenAI -> Gemini. Never hard-fails the citizen.
Q: "Is this production-ready or just a demo?" -> 90-day pilot plan with named MOEI sponsors, data residency plan, PDPL DPIA process defined""")


# ── SLIDE 8 · Arabic-First ────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, INK)
box(s, 0.15, 0, 13.18, 0.9, INK)
txt(s, "Arabic-First. Not an Afterthought.", 0.45, 0.1, 12.0, 0.75, size=30, bold=True, color=WHITE)

feats = [
    ("Jais-Family-30B", "Khaliji Arabic dialect — purpose-built for Gulf Arabic, not translated from English."),
    ("Native RTL UI", "Right-to-left is first-class: web, chat, and mobile all render RTL natively. Not a CSS patch."),
    ("Code-switching", "Arabic + English in the same turn handled cleanly — because citizens actually talk that way."),
    ("Bilingual RAG", "Postgres FTS knowledge base indexed in both languages. Ask in Arabic, get Arabic-cited results."),
]
for i, (title, desc) in enumerate(feats):
    r, c = divmod(i, 2)
    x = 0.4 + c * 6.5
    y = 1.1 + r * 2.6
    box(s, x, y, 6.2, 2.3, LGRAY)
    box(s, x, y, 0.12, 2.3, RED if c == 0 else GREEN)
    txt(s, title, x+0.3, y+0.22, 5.7, 0.6, size=22, bold=True, color=INK)
    txt(s, desc,  x+0.3, y+0.85, 5.7, 1.25, size=15, color=DGRAY)

box(s, 0.3, 6.25, 12.73, 0.85, INK)
txt(s, "\"حسن. خدمة حكومة الإمارات بذكاء يستحقه المواطن.\"",
    0.4, 6.33, 12.5, 0.7, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

notes(s, """[BACKUP — use if asked about Arabic support or localisation]

Key point: Jais is not Google Translate bolted onto an English model. It's a model specifically trained on Gulf Arabic.
The UI renders RTL as a first-class mode, not via CSS tricks.
Code-switching (mixing Arabic and English in one sentence) is handled at the router level.""")


# ── SLIDE 9 · Open Standards ──────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.15, 7.5, RED)

txt(s, "No Lock-In. No Single Point of Failure.",
    0.45, 0.15, 12.5, 0.85, size=30, bold=True, color=WHITE)
box(s, 0.4, 1.05, 12.53, 0.06, RED)

points = [
    ("LangGraph + Langfuse",
     "Apache-2.0 / MIT. No proprietary orchestration. No vendor owns the graph."),
    ("Provider-agnostic LLM cascade",
     "Groq  |  Cerebras  |  Claude (Anthropic)  |  OpenAI  |  Jais  |  Gemini — swap any provider in one config line. Graph logic unchanged."),
    ("$0 RAG on Postgres FTS",
     "Bilingual full-text search today. pgvector when scale demands. No vector-DB vendor required at pilot volume."),
    ("Graceful degradation everywhere",
     "LLM down -> fallback model. Model missing -> heuristic. Guardrail error -> safe default. The citizen never sees a hard fail."),
    ("Containerised for UAE sovereignty",
     "Deploy on Azure UAE North, Core42, or any on-prem cluster. Data residency is a deployment decision, not a rewrite."),
]
for i, (title, desc) in enumerate(points):
    y = 1.25 + i * 1.13
    box(s, 0.4, y, 0.08, 0.95, GREEN)
    txt(s, title, 0.65, y+0.03, 5.0, 0.45, size=17, bold=True, color=WHITE)
    txt(s, desc,  0.65, y+0.5,  12.0, 0.55, size=14, color=SLATE)

notes(s, """[BACKUP — use for Q&A on technology choices / vendor strategy]

Q: "Why not just use Azure OpenAI?"
A: "We can. The cascade includes OpenAI. But we're not locked to it. If MOEI requires on-prem, we swap one line of config and redeploy on Core42. The government is the customer — they set the residency rules, not us."

Q: "What if Groq shuts down or raises prices?"
A: "Automatic failover to the next provider in the cascade. Citizens never see it." """)


# ── SLIDE 10 · Impact ────────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, GREEN)
box(s, 0.15, 0, 13.18, 0.9, INK)
txt(s, "What This Means for Citizens  —  and for MOEI",
    0.45, 0.1, 12.5, 0.75, size=28, bold=True, color=WHITE)

metrics = [
    ("60%",    "self-service\ndeflection rate",   RED),
    ("< 90s",  "housing triage\n(was 5+ minutes)", GREEN),
    ("< 5s",   "first meaningful\nreply",          RED),
    ("4.5/5",  "target CSAT\nby Phase 3",          GREEN),
]
for i, (num, label, col) in enumerate(metrics):
    x = 0.5 + i * 3.1
    box(s, x, 1.05, 2.75, 2.9, LGRAY)
    box(s, x, 1.05, 2.75, 0.12, col)
    txt(s, num,   x+0.1, 1.25, 2.55, 1.3, size=44, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.1, 2.6,  2.55, 0.95, size=14, color=INK, align=PP_ALIGN.CENTER)

box(s, 0.4, 4.1, 12.53, 0.06, RGBColor(0xD0, 0xD8, 0xE0))
txt(s, "174 services  x  60% deflection  x  5-min baseline AHT  =  millions of citizen-hours returned annually",
    0.4, 4.3, 12.5, 0.65, size=16, color=INK, align=PP_ALIGN.CENTER)

box(s, 0.4, 5.1, 12.53, 0.9, INK)
txt(s, "UAE National AI Strategy 2031  |  AED 335 billion digital economy contribution target",
    0.5, 5.22, 12.2, 0.65, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txt(s, "Cost to run: < AED 1.00 per conversation  |  AED 25K/month infrastructure at pilot volume",
    0.4, 6.18, 12.5, 0.6, size=15, color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A (post-demo, punchy):
"What you just saw translates to real numbers. Sixty percent deflection. Housing triage from five minutes to ninety seconds. First reply under five seconds. At under one dirham per conversation."

"One hundred and seventy-four federal services. Millions of citizens. This is not a chatbot — this is a national infrastructure upgrade."

"Directly aligned with the UAE's AI Strategy 2031 and the AED 335 billion digital economy target."

[Move to slide 11.]""")


# ── SLIDE 11 · 90-Day Pilot ──────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
box(s, 0, 0, 0.15, 7.5, RED)
box(s, 0.15, 0, 13.18, 0.9, INK)
txt(s, "We Don't Need 3 Years.  We Need 90 Days.",
    0.45, 0.1, 12.5, 0.75, size=28, bold=True, color=WHITE)

phases = [
    ("Days 1–30", "Production Integration", RED, [
        "UAE PASS sandbox — real verified identity",
        "Azure UAE North / Core42 deployment",
        "Meta WhatsApp Cloud API under MOEI account",
        "PDPL DPIA signed by MOEI Data Protection Officer",
    ]),
    ("Days 31–60", "Shadow Mode Evaluation", RGBColor(0xC0, 0x70, 0x00), [
        "Hassan reads every real citizen turn — never replies",
        "Daily: Hassan recommendation vs. human decision",
        "Target: >= 80% agreement on 5,000-turn evaluation set",
        "Cost target: < AED 1.00 per conversation confirmed",
    ]),
    ("Days 61–90", "Limited GA Launch", GREEN, [
        "Live: housing arrears triage as default first-touch",
        "EnergyAgent + TransportAgent workers go live",
        "Citizen audit-trail UI in production",
        "60% deflection rate measured against baseline",
    ]),
]
for i, (days, title, col, bullets) in enumerate(phases):
    x = 0.4 + i * 4.27
    box(s, x, 1.05, 3.95, 5.85, LGRAY)
    box(s, x, 1.05, 3.95, 0.55, col)
    txt(s, days, x+0.18, 1.08, 3.6, 0.32, size=11, bold=True, color=WHITE)
    txt(s, title, x+0.15, 1.65, 3.65, 0.55, size=16, bold=True, color=INK)
    for j, b in enumerate(bullets):
        txt(s, "  " + b, x+0.15, 2.28 + j * 0.92, 3.65, 0.82, size=13, color=DGRAY)

txt(s, "Sponsors: MOEI Customer Happiness Centre  |  Sheikh Zayed Housing Programme  |  MOEI Digital Government / TDRA",
    0.4, 6.97, 12.5, 0.48, size=12, color=RGBColor(0x70, 0x88, 0xA0), align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A (fast but specific):
"Ninety days. Three phases. Three named partners who are already the right people at MOEI."

Phase 1: "Real UAE PASS, real data residency on Core42, real WhatsApp under MOEI's account. PDPL DPIA signed."
Phase 2: "Shadow mode — Hassan runs next to every real human agent. We measure accuracy. Target: 80% agreement."
Phase 3: "Sixty percent deflection on housing arrears. Two more services live. Citizen audit trail in production."

"Named sponsors: MOEI Customer Happiness Centre owns the CSAT. Sheikh Zayed Housing Programme provides the caseload. TDRA provides UAE PASS access."

[Close — slide 12.]""")


# ── SLIDE 12 · Close ─────────────────────────────────────────────────────────
s = slide()
box(s, 0, 0, 13.33, 7.5, INK)
box(s, 0, 0, 0.3, 7.5, RED)
box(s, 0.3, 7.42, 13.03, 0.08, GREEN)

txt(s, "حسن.", 0.5, 0.6, 12.3, 2.0, size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "خدمة حكومة الإمارات بذكاء يستحقه المواطن.",
    0.5, 2.65, 12.3, 1.1, size=28, color=GREEN, align=PP_ALIGN.CENTER)

box(s, 2.5, 3.9, 8.33, 0.06, RED)

txt(s, "Hassan.", 0.5, 4.1, 12.3, 1.1, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "UAE government service with the intelligence its citizens deserve.",
    0.5, 5.15, 12.3, 0.8, size=20, color=SLATE, align=PP_ALIGN.CENTER)

txt(s, "Team Hassan  |  MOEI x 42 Abu Dhabi  |  AgentEra 2026",
    0.5, 6.5, 12.3, 0.55, size=15, color=RGBColor(0x50, 0x70, 0x90), align=PP_ALIGN.CENTER)

notes(s, """PRESENTER A (closes slow and clear — do not rush this):

"We didn't build a chatbot."

[Pause two full seconds.]

"We built MOEI's unified digital brain. It knows the citizen. It empowers the employee. And it gives leadership real-time, explainable decision support."

[Look up from the screen. Look at the judges.]

"Hassan."

[Sit down. Let the silence hold for two seconds before Q&A starts. Do not fill it.]""")


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/alsaeed/Dev/MOEI-x-42AD-AgentEra-AI-Hackathon/Hassan_Hackathon_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
